from flask import Flask, request, render_template, redirect, url_for, send_from_directory,jsonify
import os
import shutil
from pathlib import Path
from imagegen.generate_script import VideoScriptGenerator
from imagegen.generate_image import main_download_images
from tts.generate_audio import main_generate_audio
from assembly.scripts.assembly_video import create_video, create_complete_srt, extract_topic_from_json
import re
from flask_cors import CORS
from blochchain import BlockchainSimulator
from flask_sqlalchemy import SQLAlchemy
import google.generativeai as genai
blockchain = BlockchainSimulator()
import random
###############################
import json
import re
import cv2
import torch
from transformers import AutoImageProcessor, AutoModelForImageClassification
import tempfile
from collections import defaultdict
from bs4 import BeautifulSoup
import requests
##############################

import fitz  # PyMuPDF for PDFs
import pptx  # python-pptx for PPTs
from werkzeug.utils import secure_filename

processor = AutoImageProcessor.from_pretrained("DunnBC22/vit-base-patch16-224-in21k_Human_Activity_Recognition")
model = AutoModelForImageClassification.from_pretrained("DunnBC22/vit-base-patch16-224-in21k_Human_Activity_Recognition")

gemini_api_key="AIzaSyDp4LnBQXjpjzuiOs8TDC6q9VzR66oLI8E"
app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
genai.configure(api_key='AIzaSyDp4LnBQXjpjzuiOs8TDC6q9VzR66oLI8E')
model1 = genai.GenerativeModel("gemini-2.0-flash-thinking-exp-01-21")
##############################


is_skilled=0



def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in {"pdf", "pptx"}

# Extract text from PDF (Limited for faster processing)
def extract_text_from_pdf(pdf_path, max_chars=800):
    text = ""
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text += page.get_text()
            if len(text) > max_chars:
                break
    return text[:max_chars]



# Extract text from PPTX (Limited for speed)
def extract_text_from_ppt(ppt_path, max_chars=800):
    text = ""
    prs = pptx.Presentation(ppt_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
                if len(text) > max_chars:
                    break
    return text[:max_chars]

# Generate Yes/No questions using Gemini AI
import json

def generate_yes_no_questions(text):
    try:
        prompt = f"""
            Based on the following study material, generate:
            - 5 Yes/No questions to check if the student knows the concept.

            **Return JSON format ONLY**, structured like this:
            {{
                "questions": [
                    "Question 1?",
                    "Question 2?",
                    "Question 3?"
                ]
            }}

            **Important**: 
            - Ensure the response is valid JSON.
            - Do not include any additional text or explanations.
            

            Study Material:
            {text}
            """
        

        response = model1.generate_content(prompt)
        print("Raw Gemini Response:", response.text)  # Debugging log

        if response and response.text:
            try:
                # Attempt to parse the response as JSON
                t=response.text[7:-4]
                return json.loads(t)
            except json.JSONDecodeError as e:
                print("Invalid JSON response:", t)  # Debugging log
                return {"error": f"Invalid JSON response: {str(e)}"}
        else:
            return {"error": "Empty response from Gemini AI"}
    except Exception as e:
        return {"error": f"Failed to generate questions: {str(e)}"}

# Suggest study techniques based on responses
def suggest_study_technique(correct_answers, total_questions):
    if total_questions == 0:
        return "Not enough data to suggest a technique."

    score = (correct_answers / total_questions) * 100
    is_skilled=score
    if score >= 80:
        return f"You're doing great! Maintain your knowledge with **Spaced Repetition**. you scored {score}%"
    elif score >= 50:
        return f"Use **Active Recall**: Focus on weak areas and test yourself frequently. you scored {score}%"
    else:
        return f"Try the **Feynman Technique**: Explain concepts in simple terms and reinforce learning. you scored {score}%"

    

# Route: Upload File and Generate Questions
@app.route("/upload", methods=["POST","GET"])
def upload_file():
    if request.method=="POST":
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"})

        file = request.files["file"]

        if file.filename == "" or not allowed_file(file.filename):
            return jsonify({"error": "Invalid file format. Upload PDF or PPTX."})

        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(file_path)

        # Extract text
        if filename.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        else:
            text = extract_text_from_ppt(file_path)
        print(text)

        # Generate Yes/No questions
        questions = generate_yes_no_questions(text)
        print(questions)

        if "questions" not in questions or not questions["questions"]:
            return jsonify({"error": "Failed to generate questions."})

        print("✅ Sending Questions to Frontend:", questions)  # Debugging log

        return jsonify({"questions": questions["questions"]})  # Fix JSON structure
    else:
        return render_template('perf_analyze.html')

# Route: Evaluate Yes/No Answers
@app.route("/evaluate", methods=["POST"])
def evaluate():
    if not request.is_json:
        return jsonify({"error": "Request must be JSON"}), 415

    data = request.get_json()
    questions = data.get("questions", [])
    answers = data.get("answers", [])

    if not questions or not answers:
        return jsonify({"error": "Missing questions or answers"}), 400

    correct_answers = answers.count("yes")
    total_questions = len(questions)

    # Generate study suggestion
    study_suggestion = suggest_study_technique(correct_answers, total_questions)

    # Generate fake data for percentile and class average
    percentile = random.randint(0, 100)  # Random percentile
    class_average = random.randint(50, 90)  # Random class average
    student_score = correct_answers / total_questions * 100  # Student score in percentage

    return jsonify({
        "study_suggestion": study_suggestion,
        "percentile": percentile,
        "class_average": class_average,
        "student_score": student_score
    })




##############################



# Cache to store already fetched resourcesfrom flask import Flask, request, jsonify, render_template
from bs4 import BeautifulSoup
import requests
import re
import time
import datetime
from urllib.parse import quote_plus, unquote

app = Flask(__name__)

RESOURCE_CACHE = {}
CLEANING_RULES = [
    (r'\*{2,}', ''),          # Remove bold/italic markdown
    (r'#+\s*', ''),           # Remove headings
    (r'^step \d+[:.]?\s*', ''),  # Remove step numbers
    (r'\n{2,}', '\n'),        # Remove excessive newlines
    (r'^\s+', ''),            # Remove leading whitespace
]

def fetch_learning_resources(topic):
    """Dynamically fetch learning resources with multiple fallback strategies"""
    cached = RESOURCE_CACHE.get(topic.lower())
    if cached:
        return cached

    resources = []
    search_attempts = [
        # Primary attempt with trusted domains
        {
            'query': f"best resources to learn {topic} site:youtube.com OR site:github.com OR site:docs.python.org OR site:freecodecamp.org",
            'validate': lambda url: any(d in url for d in ['youtube.com', 'github.com', 'docs.python.org', 'freecodecamp.org'])
        },
        # Secondary attempt with broader search
        {
            'query': f"best way to learn {topic} tutorial OR guide",
            'validate': lambda url: True  # Accept any URL as fallback
        }
    ]

    for attempt in search_attempts:
        try:
            search_query = quote_plus(attempt['query'])
            headers = {'User-Agent': 'Mozilla/5.0'}
            
            response = requests.get(
                f"https://www.google.com/search?q={search_query}&num=7",
                headers=headers,
                timeout=10
            )
            response.raise_for_status()

            soup = BeautifulSoup(response.text, 'html.parser')
            
            for result in soup.select('.tF2Cxc, .g, .yuRUbf')[:5]:
                try:
                    link = result.find('a')['href']
                    if link.startswith('/url?q='):
                        link = unquote(link[7:].split('&')[0])
                    
                    title = (result.find('h3') or result.find('h2')).get_text()
                    title = re.sub(r'\s+', ' ', title).strip()
                    
                    if attempt['validate'](link):
                        resources.append({
                            'title': title,
                            'url': link
                        })
                        if len(resources) >= 3:
                            break
                
                except Exception:
                    continue

            if resources:
                break

        except Exception:
            continue

    # Final fallback
    if not resources:
        resources = [
            {'title': f'Official {topic} Documentation', 'url': '#'},
            {'title': f'YouTube {topic} Tutorials', 'url': '#'},
            {'title': f'Interactive {topic} Course', 'url': '#'}
        ]

    RESOURCE_CACHE[topic.lower()] = resources
    return resources

def clean_roadmap_text(text):
    """Apply multiple cleaning rules to ensure consistent formatting"""
    for pattern, replacement in CLEANING_RULES:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE | re.MULTILINE)
    
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if line and not line.lower().startswith(('key principle', 'note:', 'tip:')):
            lines.append(line)
    
    return lines

@app.route('/generate_roadmap', methods=['POST','GET'])
def generate_roadmap():
    if request.method == 'POST':
        start_time = time.perf_counter()
        print('starting roadmap gen')
        if not request.is_json:
            return jsonify({
                'status': 'error',
                'message': 'JSON content required',
                'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
            }), 400

        data = request.get_json()
        learn = data.get('learn', '').strip()
        time_frame = data.get('time', '').strip()
        level = data.get('level', '').strip().lower()

        if not all([learn, time_frame, level]):
            return jsonify({
                'status': 'error',
                'message': 'Missing required fields',
                'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
            }), 400

        if level not in {'beginner', 'intermediate', 'advanced'}:
            return jsonify({
                'status': 'error',
                'message': 'Invalid level specified',
                'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
            }), 400

        try:
            prompt = f"""Create a concise learning roadmap for {learn} suitable for {time_frame} 
            for {level} level. Follow these rules STRICTLY:
            1. ONLY include actionable learning steps
            2. NEVER use markdown, headings, or motivational text
            3. Format EXACTLY like: "Description text | Duration: X days"
            4. Return NOTHING else except these steps
            5. it must be structured exactly like
            '''Learn core concepts | Duration: 3 days
Practice with examples | Duration: 5 days
Build mini-project | Duration: 1 week'''
            """
            
            # Replace with actual API call:
            response_text = model1.generate_content(prompt)
            print(response_text)
            steps = []
            resources = fetch_learning_resources(learn)
            
            for line in clean_roadmap_text(response_text.text):
                if not line or '|' not in line:
                    continue
                    
                desc, duration = line.split('|', 1)
                steps.append({
                    'description': desc.strip(),
                    'duration': duration.replace('Duration:', '').strip(),
                    'resources': resources
                })

            return jsonify({
                'status': 'success',
                'data': {
                    'topic': learn,
                    'timeframe': time_frame,
                    'level': level,
                    'steps': steps,
                    'resources': [r['title'] for r in resources],
                    'generated_at': datetime.datetime.now().isoformat(),
                    'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
                }
            })

        except requests.exceptions.RequestException:
            return jsonify({
                'status': 'error',
                'message': 'External service unavailable',
                'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
            }), 503

        except Exception:
            return jsonify({
                'status': 'error',
                'message': 'Roadmap generation failed',
                'execution_time': f"{(time.perf_counter() - start_time):.2f}s"
            }), 500
    else:
        return render_template('roadmap_gen.html')


    

#chat app stuff
   

##########

# @app.route('/chatbot', methods=['POST','GET'])
# def chat():
#     if request.method=='POST':
#         data = request.json
#         user_input = data.get("message")
#         genai.configure(api_key='AIzaSyDp4LnBQXjpjzuiOs8TDC6q9VzR66oLI8E')
#         if not user_input:
#             return jsonify({"error": "Message is required"}), 400

#         model1 = genai.GenerativeModel("gemini-2.0-flash-thinking-exp-01-21")

#         prompt = f"""You are a friendly AI tutor for students. Use simple language and these HTML tags: 
#         <h3>, <ul>, <ol>, <li>, <strong>, <em>, <details>, <summary>.

#         Follow this structure:

#         1. Start with: <h3>Understanding {user_input}</h3>
#         - Add <em>(1-sentence definition)</em>

#         2. Core Principles :
#         <ul>
#         detailed explaination of the concept asked in the form of points.
#         highlight the important terms
#         </ul>

#         3. If it's a math/science problem:
#         <h3>Step-by-Step Help</h3>
#         <ol>
#         <li>Restate the problem</li>
#         <li>Needed formulas/concepts</li>
#         <li>Solution steps</li>
#         <li>Check your answer</li>
#         </ol>

#         4. Always end with:
#         <details>
#         different sources the user can learn these concepts from and also add any citations if necessary.
#         </details>

#         Use <em> for examples and <strong> for important terms.
#         answer the question {user_input}
#         """
#         response = model1.generate_content(prompt)
#         bot_response = response.text[7:-4]

#         #chat_entry = ChatHistory(user_input=user_input, bot_response=bot_response)
#         #db.session.add(chat_entry)
#         #db.session.commit()

#         return jsonify({"response": bot_response})
#     else:
#         return render_template('chatbot.html')


CORS(app)
# Define paths
script_path = "resources/scripts/script.json"
images_path = "resources/images/"
audio_path = "resources/audio/"
font_path = "resources/font/font.ttf"
video_output_path = "resources/video/"
subtitle_output_path = "resources/subtitles/"

def create_or_check_folder(folder_path):
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
    else:
        if any(os.listdir(folder_path)):
            raise FileExistsError(f"Folder '{folder_path}' already exists and contains files.")
        
@app.route('/')
def home():
    return render_template('authe.html')

@app.route('/videogen', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        try:
            data = request.get_json()
            
            # Validate required fields
            required_fields = ['gem_api', 'topic', 'duration', 'key_points']
            if not all(field in data for field in required_fields):
                return jsonify({'success': False, 'error': 'Missing required fields'}), 400

            # Extract parameters
            gem_api = data['gem_api']
            topic = data['topic']
            duration = int(data['duration'])
            key_points = [word.strip() for word in data['key_points']]

            # Generate script
            generator = VideoScriptGenerator(api_key=gem_api)
            script = generator.generate_script(topic, duration, key_points)
            generator.save_script(script, script_path)

            # Generate images
            main_download_images(script_path, images_path)

            # Generate audio
            main_generate_audio(script_path, audio_path)

            # Prepare paths
            topic_cleaned = re.sub(r"[^A-Za-z0-9\s]+", " ", topic)
            topic_cleaned = re.sub(r"\s+", "_", topic_cleaned)[:100]
            sub_output_file = os.path.join(subtitle_output_path, f"{topic_cleaned}.srt")
            video_filename = f"{topic_cleaned}.mp4"
            video_file = os.path.join(video_output_path, video_filename)

            # Create subtitles
            create_complete_srt(script_path, audio_path, sub_output_file, chunk_size=10)
            print("creating video")
            # Create video
            create_video(images_path, audio_path, script_path, font_path, video_file, with_subtitles=True)

            # Copy video to final videos directory
            final_videos_dir = "videos"
            os.makedirs(final_videos_dir, exist_ok=True)
            dest_video_path = os.path.join(final_videos_dir, video_filename)
            shutil.copy(video_file, dest_video_path)

            # Cleanup all temporary files
            try:
                # Remove individual files
                for file_path in [script_path, sub_output_file]:
                    if os.path.exists(file_path):
                        os.remove(file_path)

                # Clean directories
                for dir_path in [images_path, audio_path, subtitle_output_path, video_output_path]:
                    if os.path.exists(dir_path):
                        for filename in os.listdir(dir_path):
                            file_path = os.path.join(dir_path, filename)
                            try:
                                if os.path.isfile(file_path) or os.path.islink(file_path):
                                    os.unlink(file_path)
                                elif os.path.isdir(file_path):
                                    shutil.rmtree(file_path)
                            except Exception as e:
                                app.logger.error(f'Failed to delete {file_path}: {str(e)}')

            except Exception as cleanup_error:
                app.logger.error(f"Cleanup error: {str(cleanup_error)}")

            return jsonify({
                'success': True,
                'videoPath': video_filename,
                'message': 'Video generated successfully'
            })

        except Exception as e:
            return jsonify({'success': False, 'error': str(e)}), 500

    return render_template('index.html')


@app.route('/videos/<video_name>')
def serve_video(video_name):
    return send_from_directory("videos", video_name)

@app.route('/api/videos', methods=['GET'])
def get_videos():
    try:
        videos_dir = "videos"
        if not os.path.exists(videos_dir):
            return jsonify({'success': False, 'error': 'Videos directory not found'}), 404
            
        videos = []
        for filename in os.listdir(videos_dir):
            if filename.endswith(".mp4"):
                videos.append({
                    'filename': filename,
                    'url': url_for('serve_video', video_name=filename),
                    'created_at': os.path.getctime(os.path.join(videos_dir, filename))
                })
                
        # Sort by creation time (newest first)
        videos.sort(key=lambda x: x['created_at'], reverse=True)
        
        return jsonify({
            'success': True,
            'videos': videos
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500
    
@app.route('/analyze', methods=['POST','GET'])
def analyze():
    if request.method=='POST':
        if 'video' not in request.files:
            return jsonify({'error': 'No video file provided'}), 400

        video = request.files['video']

        # Save the uploaded video to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as temp_video:
            video.save(temp_video.name)
            temp_path = temp_video.name

        try:
            result = analyze_video(temp_path)
            return jsonify(result)
        except Exception as e:
            return jsonify({'error': str(e)}), 500
        finally:
            # Clean up the temporary file
            os.unlink(temp_path)
    else:
        return render_template('analyze_it.html')

def preprocess_frame(frame):
    frame = cv2.resize(frame, (224, 224))
    frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
    inputs = processor(images=frame, return_tensors="pt")
    return inputs

def predict_action(inputs):
    with torch.no_grad():
        outputs = model(**inputs)
        predictions = outputs.logits.softmax(dim=-1)
        predicted_idx = predictions.argmax().item()
        return model.config.id2label[predicted_idx]

def analyze_video(video_path):
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise Exception("Error opening video file")

    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

    activity_durations = defaultdict(float)
    frame_count = 0
    sample_rate = 5

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        frame_count += 1
        if frame_count % sample_rate != 0:
            continue

        inputs = preprocess_frame(frame)
        action = predict_action(inputs)
        activity_durations[action] += sample_rate / fps

    cap.release()

    total_duration = sum(activity_durations.values())
    return {
        "activities": dict(activity_durations),
        "total_duration": total_duration
    }


if __name__ == "__main__":
    create_or_check_folder(images_path)
    create_or_check_folder(audio_path)
    create_or_check_folder(video_output_path)
    create_or_check_folder(subtitle_output_path)
    app.run(host='0.0.0.0',port=5000)